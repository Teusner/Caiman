# -*- coding: utf-8 -*-
"""Caiman PRe defence deck — 24 August 2026, 20 minutes, video call.

Slide text is deliberately sparse: keywords on the slide, the full script in the
speaker notes, and one bold "anchor" line per slide to fall back on.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = "C:/Users/xfalc/OneDrive/Documentos/GitHub/Caiman/rapport_stage/figures/"
ICO = os.path.join(HERE, "icons") + os.sep
import sys
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(HERE, "Caiman_PRe_soutenance.pptx")

# ---------------- palette : deep ocean ----------------
INK, DEEP, TEAL, ACCENT = "0C2231", "065A82", "1C7293", "E0A33E"
SURFACE, BODY, MUTED, WHITE = "EDF2F5", "1F3644", "5E7686", "FFFFFF"
OKC, WARNC, NOC = "2E7D5B", "B07818", "9E3B3E"
ICE, PALE, MOSSBG, ROSEBG, AMBERBG, DARKCARD = "AFC4D2", "E2EBF0", "E7F1EB", "F7EAEA", "FBF3E4", "16323F"
DEEPRED, DEEPAMBER, DEEPGREEN = "7A2A2C", "6E4A0F", "1F5641"

HEAD, SANS, MONO = "Cambria", "Calibri", "Courier New"
M, W, H = 0.6, 13.333, 7.5
CW = W - 2 * M

ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
prs.core_properties.author = "Joab da Silva Bezerra"
prs.core_properties.title = "Caiman PRe defence"
BLANK = prs.slide_layouts[6]


def rgb(h):
    return RGBColor.from_string(h)


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(bg)
    return s


def shadow(shape, blur=7, dist=1, color="9FB3C0", alpha=20):
    spPr = shape._element.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    spPr.append(etree.fromstring(
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="%d" dist="%d" dir="5400000" rotWithShape="0">'
        '<a:srgbClr val="%s"><a:alpha val="%d"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>' % (int(blur * 12700), int(dist * 12700), color, alpha * 1000)))


def shape(s, kind, x, y, w, h, fill, radius=None, shadowed=False):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = rgb(fill)
    sp.line.fill.background()
    sp.shadow.inherit = False
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    if shadowed:
        shadow(sp)
    return sp


def card(s, x, y, w, h, fill=SURFACE, radius=0.045, shadowed=True):
    adj = radius / min(w, h) if min(w, h) else 0.05
    return shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill,
                 radius=min(adj, 0.5), shadowed=shadowed)


def _style(run, d):
    f = run.font
    f.name = d.get("font", SANS)
    f.size = Pt(d.get("size", 13))
    f.bold = d.get("bold", False)
    f.italic = d.get("italic", False)
    f.color.rgb = rgb(d.get("color", BODY))
    if d.get("spc"):
        run.font._rPr.set('spc', str(int(d["spc"] * 100)))
    if d.get("sub"):
        run.font._rPr.set('baseline', '-25000')


def _bullet(p, marL=171450):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(marL))
    pPr.set('indent', str(-marL))
    bf = etree.SubElement(pPr, qn('a:buFont'))
    bf.set('typeface', 'Arial')
    bc = etree.SubElement(pPr, qn('a:buChar'))
    bc.set('char', u'\u2022')


def box(s, x, y, w, h, anchor="top", align="left"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = ANCHOR[anchor]
    tf.paragraphs[0].alignment = ALIGN[align]
    return tf


def txt(s, x, y, w, h, text, size=13, bold=False, color=BODY, font=SANS,
        align="left", anchor="top", italic=False, ls=None, spc=None, upper=False):
    tf = box(s, x, y, w, h, anchor, align)
    for i, line in enumerate((text.upper() if upper else text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN[align]
        if ls:
            p.line_spacing = ls
        r = p.add_run()
        r.text = line
        _style(r, dict(size=size, bold=bold, color=color, font=font, italic=italic, spc=spc))
    return tf


def rich(s, x, y, w, h, parts, align="left", anchor="top", ls=None):
    tf = box(s, x, y, w, h, anchor, align)
    p = tf.paragraphs[0]
    p.alignment = ALIGN[align]
    if ls:
        p.line_spacing = ls
    for d in parts:
        r = p.add_run()
        r.text = d["text"]
        _style(r, d)
        if d.get("br"):
            p = tf.add_paragraph()
            p.alignment = ALIGN[align]
            if ls:
                p.line_spacing = ls
    return tf


def bullets(s, x, y, w, h, items, size=13.5, gap=7, color=BODY, bold=False, ls=1.06):
    tf = box(s, x, y, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = ls
        p.space_after = Pt(gap)
        _bullet(p)
        r = p.add_run()
        r.text = it
        _style(r, dict(size=size, color=color, bold=bold))
    return tf


PILLC = {"ok": (OKC, WHITE), "mod": (WARNC, WHITE), "no": (NOC, WHITE),
         "live": (ACCENT, INK), "info": (TEAL, WHITE)}


def pill(s, x, y, label, kind, w=None, size=8.5):
    fill, fg = PILLC[kind]
    w = w or (0.075 * len(label) + 0.34)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.25, fill, radius=0.5)
    txt(s, x, y, w, 0.25, label, size=size, bold=True, color=fg,
        align="center", anchor="middle", spc=0.7, upper=True)
    return w


def badge(s, x, y, d, icon, fill=DEEP):
    """Icon in a filled circle — the deck's repeated motif."""
    shape(s, MSO_SHAPE.OVAL, x, y, d, d, fill)
    pad = d * 0.26
    s.shapes.add_picture(ICO + icon + ".png", Inches(x + pad / 2), Inches(y + pad / 2),
                         Inches(d - pad), Inches(d - pad))


def title(s, text, size=33, color=INK, y=0.42):
    txt(s, M, y, CW, 0.74, text, size=size, bold=True, color=color, font=HEAD, anchor="middle")


def kicker(s, text, color=TEAL, y=0.16):
    txt(s, M, y, CW, 0.26, text, size=10, bold=True, color=color, spc=1.4,
        anchor="middle", upper=True)


def caption(s, x, y, w, text, align="center", color=MUTED):
    txt(s, x, y, w, 0.3 + 0.16 * text.count("\n"), text, size=9.5, italic=True,
        color=color, align=align, ls=1.05)


def anchor_line(s, text, color=TEAL, y=6.72):
    """The one sentence to read aloud if the thread is lost."""
    txt(s, M, y, CW, 0.42, text, size=14.5, italic=True, color=color, anchor="middle")


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def logos(s):
    s.shapes.add_picture(FIG + "logos/labsticc_logo_officiel.png",
                         Inches(8.75), Inches(6.42), Inches(1.35), Inches(0.6))
    card(s, 10.35, 6.3, 2.38, 0.8, WHITE)
    s.shapes.add_picture(FIG + "logos/ensta_ip_paris_horizontal_bleu.png",
                         Inches(10.59), Inches(6.435), Inches(1.9), Inches(0.53))


def pic(s, path, x, y, w, h):
    return s.shapes.add_picture(FIG + path, Inches(x), Inches(y), Inches(w), Inches(h))


def callout(s, x, y, w, icon, head, sub, fill=DEEP):
    """Icon badge + short headline + one supporting clause."""
    badge(s, x, y, 0.62, icon, fill)
    txt(s, x + 0.82, y - 0.02, w - 0.82, 0.34, head, size=14.5, bold=True, color=DEEP, anchor="middle")
    txt(s, x + 0.82, y + 0.34, w - 0.82, 0.5, sub, size=11.5, color=BODY, ls=1.06)


def wire(s, x0, y0, x1, y1, color=None, t=0.022):
    """Orthogonal connector segment for the architecture diagram."""
    color = color or "9FB8C7"
    if abs(y1 - y0) < 1e-6:
        shape(s, MSO_SHAPE.RECTANGLE, min(x0, x1), y0 - t / 2, abs(x1 - x0), t, color)
    else:
        shape(s, MSO_SHAPE.RECTANGLE, x0 - t / 2, min(y0, y1), t, abs(y1 - y0), color)


def blk(s, x, y, w, h, header, parts, fill=PALE, hcol=DEEP):
    """Bus block: interface name over the devices carried on it."""
    card(s, x, y, w, h, fill, shadowed=False)
    txt(s, x + 0.2, y + 0.13, w - 0.4, 0.3, header, size=12.5, bold=True, color=hcol,
        anchor="middle")
    txt(s, x + 0.2, y + 0.45, w - 0.4, h - 0.55, parts, size=11, color=BODY, ls=1.12)


def dtable(s, x, y, w, cols, rows, rh=0.34, hh=0.34, size=11, hsize=10,
           mono_from=1, hot=None):
    """Data table: dark header, zebra rows, monospaced numeric columns."""
    tw = sum(c[1] for c in cols)
    k = w / tw
    shape(s, MSO_SHAPE.RECTANGLE, x, y, w, hh, INK)
    cx = x
    for i, (label, cw) in enumerate(cols):
        txt(s, cx + 0.08, y, cw * k - 0.12, hh, label, size=hsize, bold=True, color=WHITE,
            anchor="middle", align="left" if i == 0 else "center")
        cx += cw * k
    for r, row in enumerate(rows):
        ry = y + hh + r * rh
        shape(s, MSO_SHAPE.RECTANGLE, x, ry, w, rh, WHITE if r % 2 else SURFACE)
        cx = x
        for i, cell in enumerate(row):
            is_hot = hot is not None and (r, i) in hot
            txt(s, cx + 0.08, ry, cols[i][1] * k - 0.12, rh, cell,
                size=size, bold=(i == 0 or is_hot),
                font=SANS if i < mono_from else MONO,
                color=NOC if is_hot else (INK if i == 0 else BODY),
                anchor="middle", align="left" if i == 0 else "center")
            cx += cols[i][1] * k
    return y + hh + len(rows) * rh


def eqline(s, x, y, w, expr, result, size=12):
    """One worked expression with its evaluated result."""
    txt(s, x, y, w * 0.68, 0.32, expr, size=size, font=MONO, color=BODY, anchor="middle")
    txt(s, x + w * 0.68, y, w * 0.32, 0.32, result, size=size + 1, bold=True, color=DEEP,
        font=MONO, anchor="middle", align="right")


# =====================================================================
# 1 — TITLE
# =====================================================================
s = slide(INK)
txt(s, M, 1.7, 11.0, 1.0, "Design of a Subsurface Robot Swarm",
    size=42, bold=True, color=WHITE, font=HEAD, anchor="middle")
txt(s, M, 2.78, 11.0, 0.42, "Embedded electronics, simulation and prototyping of Caiman robots fleet",
    size=17, color=ACCENT, anchor="middle")
rich(s, M, 3.95, 11.0, 1.0, [
    dict(text="Joab da Silva Bezerra", size=17, bold=True, color=WHITE, br=True),
    dict(text="CNRS – Lab-STICC (UMR 6285) · ROBEX · ENSTA Brest · June to August 2026",
         size=13, color=ICE),
], ls=1.3)
txt(s, M, 5.05, 11.0, 0.6,
    "Supervisors: Franck Ruffier (CNRS) · Quentin Brateau\nAcademic tutor: Prof. Philippe Xu",
    size=12, color=MUTED, ls=1.25)
pill(s, M, 5.98, "Defence — 24 August 2026", "live", 2.35)
logos(s)
notes(s, """[0:30 - cumulative 0:30]

Good morning. I am Joab da Silva Bezerra. This work was done at the Lab-STICC on the ENSTA Brest campus, in the ROBEX team, supervised by Franck Ruffier, with technical support from Quentin Brateau.

The presentation covers the Caiman board, its architecture and production file, a mission simulator, and a hardware-in-the-loop bench.

Move on at 30 seconds.""")

# =====================================================================
# 2 — INTRODUCTION
# =====================================================================
s = slide()
kicker(s, "The system")
title(s, "Caiman: a low-cost AUV built to work in a fleet")
txt(s, M, 1.32, 7.4, 0.62,
    "A compact autonomous underwater vehicle for shallow-water survey. The unit is designed to be "
    "cheap enough to build in quantity, so a mission is flown by several vehicles rather than by "
    "one large one.", size=13, ls=1.12)
callout(s, M, 2.08, 7.4, "survey", "Bathymetric survey",
        "Mapping the bottom of a lake, a harbour or a coastal area.")
callout(s, M, 3.06, 7.4, "plume", "Plume and water-column tracking",
        "Following a pollution front or a gradient across a survey zone.")
callout(s, M, 4.04, 7.4, "inspect", "Inspection of submerged structures",
        "Close-range passes over immersed works and moorings.")

card(s, M, 5.06, 7.4, 1.72, PALE)
txt(s, M + 0.26, 5.18, 6.9, 0.3, "Why a fleet, in numbers", size=14, bold=True, color=DEEP)
rich(s, M + 0.26, 5.52, 1.9, 0.55, [
    dict(text="19.5×", size=27, bold=True, color=NOC, font=HEAD),
], anchor="middle")
txt(s, M + 2.05, 5.5, 5.1, 0.58,
    "fewer track-kilometres for the reconnaissance grid\nactually flown, at the 1.33 m sonar footprint",
    size=11, color=BODY, ls=1.12)
txt(s, M + 0.26, 6.14, 6.9, 0.56,
    "0.32 km² polygon: 240.6 km to insonify directly, 12.3 km on the 26 m grid — 2.46 km per "
    "vehicle across a fleet of five. That ratio is what makes cheap, expendable units workable.",
    size=11, color=BODY, ls=1.12)

pic(s, "fabrication/cad_caiman_full_assembly.png", 8.7, 1.45, 3.85, 4.02)
caption(s, 8.3, 5.58, 4.6,
        "Sealed PMMA cylinder · Li-ion 18650 pack · two T60 thrusters\n"
        "Bar30 depth and Ping2 sonar · nRF24 surface radio")
notes(s, """[1:55 - cumulative 2:25]

What the object is, and the number that justifies the whole approach.

A Caiman is a small AUV for shallow-water survey: a sealed PMMA cylinder, a 12 volt 56 amp-hour Li-ion pack, two T60 thrusters, a Bar30 depth sensor, a Ping2 echosounder and an nRF24 surface radio.

Three tasks it targets: bathymetric survey, mapping the bottom of a lake, harbour or coastal zone; plume and water-column tracking, following a pollution front or a gradient; and inspection of submerged structures at close range.

None of that is new on its own. What the project tests is the fleet argument, and I want to state it as a number rather than an adjective.

The Ping2 is a single-beam echosounder with a 25 degree cone. At 3 metres altitude the geometric footprint is 2 times 3 times the tangent of 12.5 degrees, so 1.33 metres. The survey polygon is 0.32 square kilometres. Insonifying every square metre of it at that footprint needs 240.6 kilometres of track. The reconnaissance grid actually flown, at 26 metre spacing, needs 12.3 kilometres.

That is a factor of 19.5. Split across five vehicles it is 2.46 kilometres each, against 48.1 kilometres each for one vehicle doing the full job.

That ratio is the engineering reason for the fleet, and it is why the electronics have to be cheap: the approach only works if you can afford five of them. The actual manufacturing economics — what five assembled boards cost, and where — are on the DFM slide.""")

# =====================================================================
# 3 — RF CONSTRAINT AND LINK BUDGET
# =====================================================================
s = slide()
kicker(s, "Communication constraint")
title(s, "No usable RF link in immersion")
callout(s, M, 1.45, 7.0, "no_radio", "RF attenuation in seawater",
        "Absorption over tens of centimetres rules out Wi-Fi and mesh networking.")
callout(s, M, 2.5, 7.0, "sub", "Store-and-forward operation",
        "Each vehicle buffers its survey to microSD and exchanges at surface contacts.")
end = dtable(s, M, 3.6, 7.0,
             [("Link parameter", 2.4), ("Value", 1.2), ("Note", 1.9)],
             [["Surface range", "250 m", "nRF24, default profile"],
              ["Raw rate", "1 Mbit/s", "physical layer"],
              ["Packet loss", "3 %", "modelled"],
              ["Telemetry period", "5 s", "per vehicle"],
              ["Fleet offered load", "256 bit/s", "5 × 32 B / 5 s"],
              ["Channel occupancy", "0.026 %", "of the 1 Mbit/s link"]],
             rh=0.33, hot={(5, 1)})
txt(s, M, end + 0.18, 7.0, 0.62,
    "The radio is not the bottleneck: the fleet offers 256 bit/s against a 1 Mbit/s channel. "
    "The binding constraint is the surface window, not throughput.",
    size=12, color=BODY, ls=1.15)

dx, dw = 8.3, 4.4
card(s, dx + 0.2, 1.55, 1.5, 0.45, SURFACE, shadowed=False)
txt(s, dx + 0.2, 1.55, 1.5, 0.45, "AUV", size=11.5, bold=True, color=DEEP,
    align="center", anchor="middle")
card(s, dx + 2.7, 1.55, 1.5, 0.45, INK, shadowed=False)
txt(s, dx + 2.7, 1.55, 1.5, 0.45, "Base", size=11.5, bold=True, color=WHITE,
    align="center", anchor="middle")
for k in range(5):
    shape(s, MSO_SHAPE.RECTANGLE, dx + 1.80 + k * 0.19, 1.755, 0.09, 0.035, TEAL)
txt(s, dx + 1.5, 1.2, 1.4, 0.3, "250 m", size=10, bold=True, color=TEAL, align="center")
shape(s, MSO_SHAPE.RECTANGLE, dx, 2.22, dw, 0.04, TEAL)
shape(s, MSO_SHAPE.RECTANGLE, dx, 2.26, dw, 3.6, "D8E7EF")
for k in range(11):
    shape(s, MSO_SHAPE.RECTANGLE, dx + 0.12 + k * 0.38, 2.72, 0.2, 0.025, "7FA8BE")
txt(s, dx + 0.12, 2.80, 1.6, 0.28, "0.5 m", size=10, bold=True, color=DEEP)
card(s, dx + 0.5, 3.75, 1.7, 0.5, WHITE, shadowed=False)
txt(s, dx + 0.5, 3.75, 1.7, 0.5, "AUV", size=11.5, bold=True, color=DEEP,
    align="center", anchor="middle")
badge(s, dx + 2.6, 3.75, 0.5, "no_radio", NOC)
txt(s, dx + 0.3, 4.55, dw - 0.6, 1.1,
    "Below 0.5 m depth the adjacency is zero. Samples are timestamped and held on microSD "
    "until the next surface contact, then merged at the rendezvous.",
    size=11.5, color=DEEP, ls=1.15)
caption(s, dx, 6.0, dw, "Link condition applied in the simulator")
notes(s, """[1:10 - cumulative 3:35]

Seawater absorbs electromagnetic waves within tens of centimetres, so Wi-Fi and the mesh networking used in aerial and ground swarms are unavailable. Acoustic modems propagate but are expensive, slow and power-hungry, which does not fit a low-cost fleet.

The operating mode that follows is store-and-forward: each vehicle buffers its survey to microSD and only exchanges when it surfaces.

Now the link budget, because it produces a result worth stating. The default profile is 250 metres range at 1 megabit per second, with 3 percent modelled packet loss and a 5 second telemetry period per vehicle. Five vehicles each sending 32 bytes every 5 seconds is 256 bits per second offered load.

Against a 1 megabit per second channel that is 0.026 percent occupancy. So the radio is nowhere near saturated. The binding constraint is not throughput, it is the surface window: how often a vehicle is above 0.5 metres and within 250 metres of a peer.

That is why the design effort went into the rendezvous logic and the frame format rather than into raw data rate.

The schematic on the right is the adjacency rule as applied in the simulator: above 0.5 metres a link exists, below it the adjacency is zero.""")

# =====================================================================
# 4 — SCOPE
# =====================================================================
s = slide()
kicker(s, "Scope")
title(s, "Planned scope and delivered scope")
items = [
    ("Schematic and routing", "check", OKC, "Corrected, reviewed, production file generated"),
    ("Production file", "check", OKC, "Gerbers, BOM, CPL — 5 boards ordered 23 June"),
    ("Embedded software", "tilde", WARNC, "Architecture and drivers on the target MCU"),
    ("Swarm control", "tilde", WARNC, "Deterministic simulator, five AUV"),
    ("Protocol bench", "check", OKC, "ESP32 and Raspberry Pi, bidirectional"),
    ("Tank and field trials", "cross", NOC, "Outside the internship window"),
]
cw2, gap = 3.86, 0.26
for i, (name, ic, col, sub) in enumerate(items):
    x = M + (i % 3) * (cw2 + gap)
    y = 1.75 + (i // 3) * 2.15
    badge(s, x, y, 0.68, ic, col)
    txt(s, x, y + 0.84, cw2, 0.6, name, size=15, bold=True, color=INK, ls=1.06)
    txt(s, x, y + 1.36, cw2, 0.55, sub, size=11.5, color=MUTED, ls=1.06)
notes(s, """[0:40 - cumulative 4:15]

Where each item of the subject ended.

Delivered: the schematic and routing corrections, the production file, and the protocol bench on two processors. Partial: the embedded software, which has the architecture and the drivers for the target microcontroller, and swarm control, which exists as a deterministic simulator rather than on vehicles. Outside the window: tank and field trials.

Keep this short. Read the six labels and move to the architecture.""")

# =====================================================================
# 5 — ARCHITECTURE
# =====================================================================
s = slide()
kicker(s, "Electronics")
title(s, "System architecture and interconnection")
txt(s, M, 1.3, 1.7, 0.5, "POWER", size=10, bold=True, spc=1.2, color=MUTED, anchor="middle")
pw_boxes = [("+BATT", "12 V"), ("LM74700", "reverse / inrush"),
            ("AP63205", "5 V rail"), ("AP63203", "3.3 V rail")]
pbw, pgap, px0 = 1.86, 0.42, 2.43
for i, (hd, sub) in enumerate(pw_boxes):
    x = px0 + i * (pbw + pgap)
    card(s, x, 1.26, pbw, 0.6, SURFACE, shadowed=False)
    txt(s, x, 1.31, pbw, 0.28, hd, size=12.5, bold=True, color=DEEP, align="center", anchor="middle")
    txt(s, x, 1.57, pbw, 0.26, sub, size=10, color=MUTED, align="center", anchor="middle")
    if i < 3:
        txt(s, x + pbw, 1.26, pgap, 0.6, "→", size=15, bold=True, color=ACCENT,
            align="center", anchor="middle")
wire(s, 10.20, 1.86, 10.20, 2.12)
wire(s, 6.65, 2.12, 10.20, 2.12)
wire(s, 6.65, 2.12, 6.65, 2.36)
card(s, 5.05, 2.36, 3.2, 2.1, INK)
txt(s, 5.15, 2.54, 3.0, 0.36, "STM32F765VITx", size=16.5, bold=True, color=WHITE,
    font=HEAD, align="center", anchor="middle")
txt(s, 5.15, 2.94, 3.0, 0.86, "Cortex-M7 at 216 MHz\nSTM32Cube and FreeRTOS\nsensor task at 10 Hz",
    size=11.5, color=ICE, align="center", ls=1.3)
pill(s, 5.6, 3.96, "116 parts · 2 layers · 94.8 cm²", "info", 2.1, size=7.5)
blk(s, M, 2.42, 3.7, 0.95, "I²C2 — inertial and environment",
    "LSM6DSO accel/gyro · LIS2MDL\nmagnetometer · LPS22HB barometer")
blk(s, M, 3.52, 3.7, 0.85, "I²C1 — external connector",
    "Bar30 depth and pressure sensor")
wire(s, 4.3, 2.90, 5.05, 2.90)
wire(s, 4.3, 3.95, 5.05, 3.95)
blk(s, 9.0, 2.36, 3.7, 0.78, "SPI1 — radio", "nRF24L01+ with CE, CSN and IRQ")
blk(s, 9.0, 3.24, 3.7, 0.78, "SDMMC — storage", "microSD with card-detect contact")
blk(s, 9.0, 4.12, 3.7, 0.78, "UART — navigation aids", "GNSS receiver · Ping2 echosounder")
for yy in (2.75, 3.63, 4.51):
    wire(s, 8.25, yy, 9.0, yy)
wire(s, 6.65, 4.46, 6.65, 5.0)
wire(s, 2.53, 5.0, 10.77, 5.0)
bot = [("PWM — propulsion", "Two ESC channels driving the T60 thrusters"),
       ("GPIO — signalling", "SK6812 RGB via SN74LV1T34 level shifter"),
       ("ADC and GPIO — housekeeping", "Battery · US1881 magnetic contact · leak")]
bbw = 3.86
for i, (hd, sub) in enumerate(bot):
    x = M + i * (bbw + 0.26)
    wire(s, x + bbw / 2, 5.0, x + bbw / 2, 5.16)
    blk(s, x, 5.16, bbw, 0.85, hd, sub, SURFACE)
dtable(s, M, 6.2, 12.1,
       [("Buck converter", 2.0), ("L", 1.0), ("D", 0.8), ("f_sw", 1.0), ("ΔI_L", 1.0),
        ("ΔI / I_out", 1.1), ("I_pk at 2 A", 1.2), ("I_sat", 1.0), ("Margin", 1.0)],
       [["L1 — AP63205, 5 V", "4.7 µH", "0.42", "1.1 MHz", "0.56 A", "28.2 %", "2.28 A", "2.5 A", "8.7 %"],
        ["L2 — AP63203, 3.3 V", "3.3 µH", "0.66", "1.1 MHz", "0.31 A", "15.5 %", "2.16 A", "3.0 A", "28.2 %"]],
       rh=0.3, hh=0.3, size=10.5, hsize=9.5, hot={(0, 8)})
notes(s, """[2:00 - cumulative 6:15]

The board is organised around one microcontroller, an STM32F765VITx, Cortex-M7 at 216 megahertz, running STM32Cube and FreeRTOS with the sensor task at 10 hertz. 116 placed components on two layers, 94.8 square centimetres.

Power enters at the top. The 12 volt input is protected by an LM74700 ideal-diode controller against reverse polarity and inrush, then two cascaded synchronous bucks: an AP63205 for the 5 volt rail and an AP63203 deriving 3.3 volts from it.

The interfaces are grouped by function. On the left, two I2C buses, deliberately separated. I2C2 is the internal instrumentation bus carrying the LSM6DSO accelerometer and gyroscope, the LIS2MDL magnetometer and the LPS22HB barometer. I2C1 is brought out to a connector for the Bar30, which sits in the front end cap and therefore has to leave the board.

On the right, SPI1 for the nRF24 transceiver with its chip-enable, chip-select and interrupt lines; SDMMC for the microSD with card-detect, which is what makes store-and-forward possible; and a UART group for the GNSS receiver and the Ping2 echosounder.

At the bottom, two PWM channels for the ESCs driving the T60 thrusters, a GPIO line driving the SK6812 indicator through an SN74LV1T34 level shifter because the LED needs 5 volt logic, and an ADC and GPIO group for battery voltage, the US1881 Hall-effect contact and the leak detector. The magnetic contact exists because the hull is sealed: you cannot fit a switch.

Now the converter table at the bottom, which is the quantitative part.

Applying the standard buck relation, delta-I equals V-out times one minus D, over L times switching frequency. On L1 that gives 0.56 amps of ripple, which is 28.2 percent of the 2 amp rating. On L2, 0.31 amps, 15.5 percent. Both sit inside the usual 20 to 40 percent design band, L2 conservatively so.

Peak current is output plus half the ripple. L1 reaches 2.28 amps against a rated saturation of 2.5, so 8.7 percent margin. L2 reaches 2.16 against 3.0 amps, 28.2 percent margin. L1 is the tighter of the two and is the rail I would instrument first.""")

# =====================================================================
# 6 — FIRMWARE
# =====================================================================
s = slide()
kicker(s, "Embedded software")
title(s, "Firmware structure and task budget")
end = dtable(s, M, 1.45, 6.9,
             [("Module group", 2.5), ("Devices", 2.7), ("State", 1.4)],
             [["Inertial and environment", "LSM6DSO, LIS2MDL, LPS22HB", "in main loop"],
              ["Storage and housekeeping", "microSD, battery ADC, RGB", "in main loop"],
              ["Radio, thrusters, depth", "nRF24, ESC, Bar30", "driver written"],
              ["Navigation aids and safety", "GNSS, Ping2, RC, leak, contact", "skeleton"]],
             rh=0.50)
txt(s, M, end + 0.24, 6.9, 0.9,
    "STM32Cube HAL with FreeRTOS. The application task initialises the sensor layer, then calls "
    "its service function on a 100 ms period, so a 10 Hz acquisition cadence against a 5 s "
    "telemetry period — fifty acquisition cycles per transmitted frame.",
    size=12.5, color=BODY, ls=1.15)
card(s, 7.75, 1.45, 4.95, 2.55, PALE)
txt(s, 8.0, 1.6, 4.45, 0.34, "Timing budget", size=14, bold=True, color=DEEP)
eqline(s, 8.0, 2.06, 4.45, "sensor task period", "100 ms")
eqline(s, 8.0, 2.50, 4.45, "acquisition rate", "10 Hz")
eqline(s, 8.0, 2.94, 4.45, "telemetry period", "5 s")
eqline(s, 8.0, 3.38, 4.45, "cycles per frame", "50")
card(s, 7.75, 4.25, 4.95, 1.7, AMBERBG)
txt(s, 8.0, 4.42, 4.45, 1.36,
    "Logger defect to correct: the CSV header declares ten columns while the application writes "
    "eighteen fields. It has to be reconciled before the logs are used quantitatively.",
    size=12, color=DEEPAMBER, ls=1.15)
notes(s, """[1:00 - cumulative 7:15]

The firmware is STM32Cube HAL with FreeRTOS. The application task initialises the sensor layer and then calls its service function on a 100 millisecond period, so 10 hertz acquisition.

Put that against the 5 second telemetry period from the link budget and you get 50 acquisition cycles per transmitted frame. Which is the reason the frame carries a compacted snapshot rather than a sample stream: at 10 hertz you cannot transmit every sample through a surface-only link, so the vehicle logs locally and sends state.

By module group: the inertial and environment sensors and the storage and housekeeping group are called from the main loop. The radio, the ESCs and the Bar30 have drivers written but sit outside the audited application path. The navigation aids and the safety inputs are skeletons.

One defect I will state because it affects data quality: the CSV logger declares ten columns in its header and writes eighteen fields. That has to be reconciled before the logs are used quantitatively.""")

# =====================================================================
# 7 — SCHEMATIC CORRECTION
# =====================================================================
s = slide()
kicker(s, "Electronics review")
title(s, "Schematic corrections on the inherited design")
callout(s, M, 1.45, 7.0, "capacitor", "Decoupling placed in series",
        "C12–C18 sat in series on V_DD, V_DDA and V_BAT instead of in parallel to ground.")
card(s, M, 2.55, 7.0, 0.95, ROSEBG)
txt(s, M + 0.3, 2.65, 6.4, 0.75,
    "A series capacitor is an open circuit at DC: the 3.3 V rail did not reach pins "
    "11, 27, 50, 75, 100 and 71.",
    size=14, bold=True, color=DEEPRED, anchor="middle", ls=1.12)
end = dtable(s, M, 3.68, 7.0,
             [("Commit", 1.3), ("Correction", 3.6), ("Effect", 1.6)],
             [["dd041de", "C12–C18 returned to parallel, rail to GND", "DC continuity"],
              ["c66b3fb", "LSM6DSO ground pin electrical type", "matrix conflict cleared"],
              ["2be6335", "Redundant power flags removed", "net declarations clean"]],
             rh=0.4)
card(s, M, end + 0.2, 7.0, 0.8, INK)
rich(s, M + 0.3, end + 0.3, 6.4, 0.6, [
    dict(text="10 → 0", size=25, bold=True, color=ACCENT, font=HEAD),
    dict(text="   blocking ERC errors, KiCad 10.0.5 (42 minor warnings remain)",
         size=13, color=WHITE),
], anchor="middle")
pic(s, "fabrication/stm32_decoupling_capacitors_before.png", 8.05, 1.45, 2.4, 2.82)
pic(s, "fabrication/stm32_decoupling_capacitors_fix.png", 10.68, 1.45, 2.0, 2.82)
txt(s, 8.05, 4.35, 2.4, 0.26, "BEFORE", size=11, bold=True, color=NOC, align="center", spc=0.9)
txt(s, 10.68, 4.35, 2.0, 0.26, "AFTER", size=11, bold=True, color=OKC, align="center", spc=0.9)
caption(s, 8.05, 4.68, 4.63, "Series on the rail, then parallel to ground")
notes(s, """[1:15 - cumulative 8:30]

The corrections applied to the inherited schematic.

The significant one: the 100 nanofarad decoupling capacitors on VDD, VDDA and VBAT were drawn in series on the supply rails rather than in parallel to ground. A capacitor in series is an open circuit at DC, so the 3.3 volt rail did not reach pins 11, 27, 50, 75, 100 and 71.

Commit dd041de returns C12 to C18 to parallel between the rail and ground, restoring DC continuity while keeping the high-frequency decoupling the reference manual specifies.

Two further corrections in the table: c66b3fb fixes the electrical type of the LSM6DSO ground pin, which was producing a conflict in the KiCad interconnection matrix, and 2be6335 removes redundant power flags.

The measurable outcome is blocking ERC errors going from ten to zero under KiCad 10.0.5. Forty-two minor warnings remain and are listed in the report.""")

# =====================================================================
# 8 — DRC
# =====================================================================
s = slide()
kicker(s, "Routing")
title(s, "Design rules and manufacturing margins")
end = dtable(s, M, 1.45, 11.9,
             [("Constraint", 2.3), ("Project rule", 1.4), ("Process capability", 1.7),
              ("Margin", 1.1)],
             [["Internal copper clearance", "0.150 mm", "0.100 mm", "+50 %"],
              ["Minimum track width", "0.127 mm", "0.100 mm", "+27 %"],
              ["Minimum via diameter", "0.300 mm", "0.200 mm", "+50 %"]],
             rh=0.36)
txt(s, M, end + 0.2, 11.9, 0.4,
    "Rules were set from the published tolerance grid and then tightened, so the layout keeps "
    "margin against the process floor rather than sitting on it.",
    size=12, color=BODY, ls=1.15)
end2 = dtable(s, M, end + 0.7, 11.9,
              [("Documented exclusion", 2.6), ("Measured", 1.3), ("Against capability", 1.8)],
              [["U5 — /I2C2_SCL to GND", "0.1291 mm", "+29 % above 0.100 mm"],
               ["U8 — INT2 pad", "0.1472 mm", "+47 % above 0.100 mm"],
               ["U2 — BST / SW pins", "0.1320 mm", "vendor footprint, fixed"],
               ["U6 — unconnected pad", "—", "no electrical function"]],
              rh=0.34)
txt(s, M, end2 + 0.18, 11.9, 0.4,
    "Every deviation sits above the published process capability and is recorded in the project "
    "file with its justification, rather than removed by relaxing the rule set globally.",
    size=12, color=BODY, ls=1.15)
card(s, M, end2 + 0.68, 11.9, 0.85, INK)
rich(s, M + 0.28, end2 + 0.78, 11.34, 0.65, [
    dict(text="Lesson.  ", size=13, bold=True, color=ACCENT),
    dict(text="DRC belongs at the exit of the production export, re-run on the exact state that "
              "feeds the transmitted files, not a milestone cleared once mid-routing.", size=13,
         color=WHITE),
], anchor="middle", ls=1.15)
notes(s, """[1:05 - cumulative 9:35]

The design rules and the margin they leave.

The top table: the project rules were taken from the manufacturer's published tolerance grid and then tightened. Internal copper clearance 0.150 millimetres against a 0.100 capability, so 50 percent margin. Minimum track 0.127 against 0.100, 27 percent. Minimum via 0.300 against 0.200, 50 percent.

The second table lists the four deviations that were accepted, with their measured values. The largest, U5, is 0.1291 millimetres between I2C2_SCL and ground inside the package, 29 percent above the 0.100 millimetre capability. U8's INT2 pad is 0.1472, so 47 percent above capability. U2's BST and SW pins are set by the manufacturer's footprint. U6 is an unconnected pad with no electrical function.

Each of these sits above what the process can actually hold, and each is recorded in the project file with its justification rather than hidden by lowering the rule set globally.

The lesson on the dark card is the one I keep repeating because it is the real methodological outcome of this phase: DRC belongs at the exit of the production export, re-run on the exact state that feeds the transmitted files, not a milestone you clear once during routing.

>>> IF ASKED whether it was re-run on the transmitted revision: Q&A sheet, answer 1, five points in order. Do not improvise.

Next slide is the actual routed copper.""")

# =====================================================================
# 9 — ROUTING
# =====================================================================
s = slide()
kicker(s, "Routing")
title(s, "Routed copper: top and bottom layers")
img_h, img_gap = 4.05, 0.26
img_w = img_h * (2143 / 2208)
total_w = img_w * 2 + img_gap
x0 = (W - total_w) / 2
y0 = 1.5
s.shapes.add_picture(ICO + "routing_top.png",
                     Inches(x0), Inches(y0), Inches(img_w), Inches(img_h))
s.shapes.add_picture(ICO + "routing_bottom.png",
                     Inches(x0 + img_w + img_gap), Inches(y0), Inches(img_w), Inches(img_h))
caption(s, x0, y0 + img_h + 0.08, img_w, "Top layer (F.Cu) — power stage, MCU, SMD footprints")
caption(s, x0 + img_w + img_gap, y0 + img_h + 0.08, img_w,
        "Bottom layer (B.Cu) — ground plane and decoupling")
txt(s, M, 6.28, CW, 0.32,
    "100.87 × 93.94 mm · 2 layers · 1 oz Cu · ENIG · black mask, white silkscreen · tented vias "
    "· IPC class 2 · flying-probe tested",
    size=11.5, color=MUTED, align="center", anchor="middle")
notes(s, """[0:55 - cumulative 10:30]

The routed board itself, both copper layers, from the same KiCad revision the production file was exported from.

Top layer on the left: the power stage top centre, buck converters and their inductor footprints, the STM32 fanout, and the SMD footprints for the passives. Bottom layer on the right: the ground plane pour with stitching vias, and the decoupling network under the sensor and radio ICs.

Spec at the bottom: 100.87 by 93.94 millimetres, two layers, 1 ounce copper, ENIG finish, black solder mask with white silkscreen, tented vias, IPC class 2, flying-probe tested — the same figures cited on the manufacturing order.

COMPRESSIBLE: this slide is a visual pause, not new claims. Move quickly if behind.""")

# =====================================================================
# 10 — COMPONENT SELECTION
# =====================================================================
s = slide()
kicker(s, "Design for manufacturability")
title(s, "Power inductor selection")
pic(s, "fabrication/jlcpcb_pcba_L1L2_pads.png", M, 1.5, 4.3, 3.38)
caption(s, M, 4.98, 4.3, "Board pads in green against the candidate terminations in red")
txt(s, 5.35, 1.45, 7.35, 0.62,
    "The PCBA engineering review flagged that the pre-selected inductor terminations exceeded the "
    "drawn land pattern. Selection was redone against pad geometry as a first-order constraint.",
    size=13, color=BODY, ls=1.15)
end = dtable(s, 5.35, 2.2, 7.35,
             [("Designator", 1.3), ("Value", 0.9), ("Part", 2.0), ("LCSC", 1.2),
              ("I_sat", 0.9), ("Footprint", 1.5)],
             [["L1 — 5 V buck", "4.7 µH", "Würth 74438356047", "C2045367", "2.5 A", "WE-MAPI 4020"],
              ["L2 — 3.3 V buck", "3.3 µH", "Würth 74438356033", "C2045382", "3.0 A", "WE-MAPI 4020"]],
             rh=0.36)
txt(s, 5.35, end + 0.22, 7.35, 0.86,
    "L2 moves from the schematic value of 3.9 µH to 3.3 µH. Ripple rises from 0.26 A to "
    "0.31 A, about 18 %, and the peak stays at 72 % of saturation — the substitution costs "
    "nothing in margin.",
    size=12, color=BODY, ls=1.15)
card(s, 5.35, end + 1.1, 7.35, 0.95, INK)
txt(s, 5.65, end + 1.22, 6.75, 0.75,
    "Identical 4 × 4 mm package outlines carried incompatible pad geometry. Selection is "
    "constrained by the land pattern before it is constrained by the value.",
    size=13, bold=True, color=ACCENT, anchor="middle", ls=1.15)

card(s, M, 5.5, 12.1, 1.6, PALE)
txt(s, M + 0.26, 5.6, 8.0, 0.28,
    "Manufacturing economics — JLCPCB, Shenzhen, China (jlcpcb.com)",
    size=13, bold=True, color=DEEP)
eqline(s, M + 0.26, 5.94, 5.65, "bare PCB, 5 boards", "€18.88")
eqline(s, M + 0.26, 6.26, 5.65, "PCBA, assembled, 5 boards", "€368.87")
eqline(s, M + 0.26, 6.58, 5.65, "shipping — DHL Express, 0.83 kg", "€20.11")
shape(s, MSO_SHAPE.RECTANGLE, M + 6.15, 5.6, 0.018, 1.3, "C9D6DF")
rich(s, M + 6.55, 5.66, 2.5, 0.62, [
    dict(text="€407.87", size=26, bold=True, color=DEEP, font=HEAD),
], anchor="middle")
txt(s, M + 6.55, 6.26, 2.5, 0.3, "order total, 5 boards", size=10.5, color=MUTED)
card(s, M + 9.25, 5.66, 2.55, 1.3, WHITE, shadowed=False)
rich(s, M + 9.45, 5.78, 2.15, 0.5, [
    dict(text="€81.57", size=22, bold=True, color=NOC, font=HEAD),
], anchor="middle")
txt(s, M + 9.45, 6.26, 2.15, 0.62, "per assembled unit, landed\n(n = 5 prototype run)",
    size=10, color=BODY, ls=1.1)
notes(s, """[1:20 - cumulative 11:50]

The power inductor selection, which was redone during PCBA preparation.

The manufacturer's engineering review flagged that the terminations of the pre-selected parts exceeded the land pattern drawn on the board. You can see the overhang in the image: board pads in green, candidate terminations in red.

The selection was redone with pad geometry as a first-order constraint, and checked against inductance, rated thermal current, saturation current and DC resistance. Both positions land on the WE-MAPI 4020 footprint, L1 at 4.7 microhenry with 2.5 amps saturation, L2 at 3.3 microhenry with 3.0 amps.

The general point on the dark card: two components with identical 4 by 4 millimetre package outlines had incompatible pad geometry. Component selection is constrained by the land pattern before it is constrained by the electrical value.

And the manufacturing economics, since I have the actual invoice. This order, W2026062315153820, was fabricated and assembled by JLCPCB in Shenzhen: 18 euros 88 for five bare boards, 368 euros 87 for five assembled boards — that figure covers component procurement and SMT assembly labour — and 20 euros 11 for DHL Express shipping of the 0.83 kilogram package. Order total 407 euros 87.

Divided across five fully assembled, populated boards that is 81 euros 57 per unit, landed. This is prototype-quantity pricing at five pieces; per-unit cost falls substantially at production volumes, so I would not extrapolate this figure to a series run. But it is the real number for this batch, and it is what makes the fleet economically credible at this stage.""")

# =====================================================================
# 11 — METHOD
# =====================================================================
s = slide(INK)
kicker(s, "Method", color=ACCENT)
txt(s, M, 1.3, 12.1, 1.0, "Validating the fleet layer before the hardware arrives",
    size=34, bold=True, color=WHITE, font=HEAD, anchor="middle")
txt(s, M, 2.35, 11.4, 0.8,
    "Board fabrication and assembly run on their own lead time. The fleet behaviour, the "
    "protocol and its security properties do not depend on that board, so they were specified, "
    "implemented and exercised in parallel.",
    size=15, color=ICE, ls=1.25)
path = [("monitor", "Simulator", "Fleet logic, rendezvous, faults", TEAL),
        ("packet", "Protocol", "32-byte frame, AEAD, anti-replay", DEEP),
        ("two_chips", "HIL bench", "Same C source, two processors", OKC)]
pw = 3.6
for i, (ic, hd, sub, col) in enumerate(path):
    x = M + i * (pw + 0.65)
    card(s, x, 3.55, pw, 2.05, DARKCARD)
    badge(s, x + (pw - 0.86) / 2, 3.8, 0.86, ic, col)
    txt(s, x, 4.8, pw, 0.36, hd, size=19, bold=True, color=WHITE, font=HEAD, align="center")
    txt(s, x, 5.2, pw, 0.36, sub, size=12, color=ICE, align="center", ls=1.06)
    if i < 2:
        txt(s, x + pw + 0.06, 3.55, 0.53, 2.05, "→", size=28, bold=True, color=ACCENT,
            align="center", anchor="middle")
notes(s, """[0:55 - cumulative 12:45]

This slide states the working method for the second half of the internship.

Board fabrication and assembly run on their own lead time. But the fleet layer does not depend on that board: the mission logic, the rendezvous behaviour, the frame format and its security properties are all specifiable and testable independently.

So they were developed in parallel, in three steps. A deterministic simulator for the fleet logic, the rendezvous and the fault scenarios. A protocol specified to the physical constraint of the radio, with authenticated encryption and anti-replay. And then a hardware-in-the-loop bench running that same C source on two physical processors.

The three next slides follow that order.""")

# =====================================================================
# 12 — SIMULATOR
# =====================================================================
s = slide()
kicker(s, "Simulation")
title(s, "Deterministic mission model")
pic(s, "simulator/mission_after_first_sync.png", M, 1.45, 7.0, 4.59)
caption(s, M, 6.12, 7.0, "Mission view after the first rendezvous")
end = dtable(s, 7.75, 1.45, 4.95,
             [("Parameter", 2.2), ("Value", 1.5)],
             [["Vehicles", "5"],
              ["Compute zone", "1000 × 700 m"],
              ["Survey polygon", "0.32 km²"],
              ["Bottom depth", "13–20 m"],
              ["Survey altitude", "3 m"],
              ["Beam width", "25°"],
              ["Footprint", "1.33 m"],
              ["Grid spacing", "26 m"],
              ["Track length", "12.3 km"]], rh=0.32)
txt(s, 7.75, end + 0.2, 4.95, 0.62,
    "A fixed seed reproduces terrain, mission decisions and packet outcomes, so two runs differ "
    "only where a parameter was deliberately changed.",
    size=12, color=BODY, ls=1.15)
card(s, 7.75, end + 1.05, 4.95, 1.05, AMBERBG)
txt(s, 8.0, end + 1.16, 4.45, 0.85,
    "The station displays received telemetry only. Delayed samples are replayed from their "
    "timestamps rather than interpolated on arrival.",
    size=12, color=DEEPAMBER, anchor="middle", ls=1.15)
notes(s, """[1:40 - cumulative 14:25]

The simulator is a system-level model, not a hydrodynamic one. Its purpose is the fleet layer: mission distribution, the rendezvous, the link condition and the fault scenarios.

The parameter table on the right is the reference configuration and it ties back to the coverage budget from the second slide. Five vehicles over a thousand by seven hundred metre zone, a 0.32 square kilometre survey polygon, synthetic bottom between 13 and 20 metres. Vehicles fly boustrophedon transects at 3 metres altitude with the 25 degree beam, giving the 1.33 metre footprint, on a 26 metre grid, for the 12.3 kilometres of track we computed earlier.

Determinism: a fixed seed reproduces the terrain, the mission decisions and the packet outcomes. Two runs differ only where a parameter was deliberately changed, which is what makes a comparison between configurations meaningful.

And the design point I would highlight, on the amber card: the supervision view shows received telemetry only. After a dive the station holds the last position received by radio and propagates the commanded route with growing uncertainty; it does not read the model's internal state. When delayed telemetry arrives at a rendezvous the samples are replayed from their timestamps rather than interpolated, so the reconstructed track is the one actually flown.

That separation is what makes the supervision view a usable operator tool rather than a display of ground truth.""")

# =====================================================================
# 13 — DEMONSTRATION
# =====================================================================
s = slide()
kicker(s, "Demonstration")
title(s, "Supervision interface")
pic(s, "simulator/surface_mesh_rendezvous.png", M, 1.5, 7.6, 4.99)
pill(s, 8.4, 1.55, "Run live — 60 s", "live", 2.0)
txt(s, 8.4, 2.02, 4.3, 0.28, "Seven supervision views", size=11, bold=True,
    color=MUTED, spc=0.8, upper=True)
txt(s, 8.4, 2.42, 4.3, 2.3,
    "Mission\nRobots\nNetwork\nCommands\nSecurity\nPacket Log\nEvents",
    size=14, bold=True, color=DEEP, ls=1.5)
card(s, 8.4, 4.86, 4.3, 1.4, PALE)
txt(s, 8.66, 4.98, 3.8, 1.16,
    "Python 3.12 and Streamlit containerised behind a Caddy reverse proxy with a health "
    "endpoint. Run here from a local instance.",
    size=11.5, color=BODY, anchor="middle", ls=1.15)
notes(s, """[1:00 - cumulative 15:25]  << HARD CAP: 60 SECONDS.

>>> LIVE. Started before the call opens. Do not launch it here.

SETUP BEFORE THE CALL
  cd simulator && streamlit run app.py
  Browser at http://localhost:8501, window sized for share, default seed, five AUV.
  Advance the clock just past the first rendezvous so the grid is partly synchronised.
  Deck on a second display.

SEQUENCE - four steps, about 15 seconds each. Do not explore.
  1. Mission view. State that this shows received telemetry only. Point at one last-known marker.
  2. Advance the clock. Submerged vehicles stop updating: the adjacency rule, live.
  3. Network view at a rendezvous. Edges appear; trace the multi-hop route to the base.
  4. Security view. Tamper the AEAD tag once and show the rejection counter increment by category.

  Then stop. Do not open Packet Log or Events unless asked.

IF IT FAILS: 15 seconds maximum, then move to the backup slide.

Say once: containerised behind a Caddy reverse proxy with a health endpoint, deployed to a VPS
under the caimansim.fr domain; running here from a local instance.""")

# =====================================================================
# 14 — BACKUP
# =====================================================================
s = slide()
kicker(s, "Backup", color=NOC)
title(s, "Interface captures")
pic(s, "simulator/surface_mesh_rendezvous.png", M, 1.5, 5.9, 3.87)
caption(s, M, 5.46, 5.9, "Rendezvous: four AUV surfaced, active links, multi-hop route to base",
        align="left")
pic(s, "simulator/security_physical_capture.png", 6.8, 1.5, 5.9, 3.87)
caption(s, 6.8, 5.46, 5.9, "Capture scenario: R1 taken, key 1 compromised, renewal pending",
        align="left")
card(s, M, 5.9, 12.1, 1.16)
rich(s, M + 0.3, 6.0, 11.5, 0.96, [
    dict(text="Capture scenario.  ", size=13, bold=True, color=DEEP),
    dict(text="The compromised vehicle's key is revoked. Submerged peers receive no immediate "
              "recall: they continue to their next surface contact, take the authenticated return "
              "order there, and are rekeyed with the captured node excluded. The group key keeps "
              "the demonstration legible; a production system needs per-vehicle roots and pairwise "
              "traffic keys.", size=12.5),
], anchor="middle", ls=1.12)
notes(s, """[Skip unless the live instance failed. Budget 0:00.]

Left: network view at a rendezvous. Right: the physical capture scenario.

On the capture case: the key of the captured vehicle is revoked. The submerged peers receive no immediate recall, because there is no link to them. They continue to their next surface contact, take the authenticated return order there, and are rekeyed with the captured node excluded.

The group key keeps the demonstration legible. A production system needs per-vehicle roots and pairwise traffic keys.

Then continue to the frame format.""")

# =====================================================================
# 15 — FRAME FORMAT
# =====================================================================
s = slide()
kicker(s, "Protocol")
title(s, "Frame format and link budget")
barY, barH = 1.6, 1.6
segs = [(CW * 0.25, DEEP, "8 B", "header", WHITE, "6FA5C0"),
        (CW * 0.25, TEAL, "8 B", "encrypted payload", WHITE, "7FB4CA"),
        (CW * 0.50, ACCENT, "16 B", "Poly1305 tag", INK, "B27F22")]
bx = M
for w_, fill, big, lbl, fg, tick in segs:
    shape(s, MSO_SHAPE.RECTANGLE, bx, barY, w_, barH, fill)
    n = int(round(w_ / (CW / 32.0)))
    for k in range(1, n):
        shape(s, MSO_SHAPE.RECTANGLE, bx + k * (w_ / n) - 0.006, barY, 0.012, 0.2, tick)
    txt(s, bx, barY + 0.36, w_, 0.56, big, size=36, bold=True, color=fg, font=HEAD,
        align="center", anchor="middle")
    txt(s, bx, barY + 0.95, w_, 0.38, lbl, size=14, color=fg, align="center", anchor="middle")
    bx += w_
txt(s, M, 3.28, CW, 0.28, "one tick per byte · 32 B is the nRF24 physical payload limit",
    size=10, italic=True, color=MUTED, align="center")
end = dtable(s, M, 3.68, 6.5,
             [("Quantity", 2.6), ("Value", 1.3), ("Derivation", 2.0)],
             [["Application efficiency", "25 %", "8 B of 32 B"],
              ["Max fragments", "16", "4-bit index"],
              ["Max logical payload", "128 B", "16 × 8 B"],
              ["Frame air time", "256 µs", "256 bit at 1 Mbit/s"],
              ["With PHY overhead", "320 µs", "+8 B preamble, addr, CRC"]],
             rh=0.34)
txt(s, M, end + 0.2, 6.5, 0.62,
    "The 8 encrypted bytes carry x/y position, vehicle depth, bottom depth, battery, link "
    "quality, leak flag and GNSS availability, packed into 64 bits.",
    size=12, color=BODY, ls=1.15)
card(s, 7.15, 3.68, 5.55, 2.1, PALE)
txt(s, 7.42, 3.82, 5.0, 0.34, "Cryptographic construction", size=14, bold=True, color=DEEP)
bullets(s, 7.42, 4.22, 5.0, 1.45, [
    "256-bit per-mission master key; HKDF-SHA-256 derives separate encryption, routing and "
    "fingerprint material.",
    "ChaCha20-Poly1305 AEAD — no AES accelerator on this MCU.",
    "Monotonic per-source sequence in the nonce, anti-replay window, relay de-duplication cache.",
], size=11.5, gap=4)
txt(s, 7.15, 5.95, 5.55, 0.6,
    "Relays retransmit the 32 bytes unchanged: TTL and hop count stay in local policy, "
    "ciphertext and tag are never rewritten.",
    size=12, color=BODY, ls=1.15)
notes(s, """[1:35 - cumulative 17:00]

The nRF24 has a fixed 32-byte physical payload, so the frame was specified to that limit rather than around it.

The allocation is eight bytes of header, eight bytes of encrypted payload and sixteen bytes of Poly1305 tag. Into those 64 payload bits go x and y position, vehicle depth, bottom depth, battery, link quality, a leak flag and GNSS availability.

The table gives the derived quantities. Application efficiency is 25 percent, eight useful bytes in thirty-two. The fragment index is four bits, so sixteen fragments maximum, which bounds the logical payload at 128 bytes. Frame air time is 256 bits at 1 megabit per second, so 256 microseconds; adding the nRF24 preamble, address and CRC takes it to about 320.

Put that against the link budget from earlier: 256 microseconds of air time every 5 seconds per vehicle is why fleet occupancy came out at 0.026 percent.

On the crypto: a 256-bit per-mission master key, with HKDF-SHA-256 deriving separate material for encryption, routing and the displayed fingerprint. ChaCha20-Poly1305 for the AEAD, chosen because this microcontroller has no AES accelerator and ChaCha is efficient and constant-time in software. A monotonic per-source sequence goes into the nonce, with an anti-replay window and a de-duplication cache at relays.

And relays retransmit the 32 bytes unchanged: TTL and hop count are local policy, the ciphertext and its tag are never rewritten.

IF ASKED about 25 percent: truncating the tag trades forgery resistance for payload, and authenticating a fragment group rather than each fragment would amortise it. The full tag was chosen because the objective was the security property, not throughput — and at 0.026 percent occupancy there is no throughput pressure to trade against.""")

# =====================================================================
# 16 — HIL
# =====================================================================
s = slide()
kicker(s, "Hardware-in-the-loop")
title(s, "Protocol running on two processors")
pic(s, "hil/hil_banc_2026-08-16.jpg", M, 1.45, 6.3, 4.37)
caption(s, M, 5.92, 6.3, "16 August 2026: ESP32 as R1, Raspberry Pi 4 as R2")
end = dtable(s, 7.2, 1.45, 5.5,
             [("Node", 1.2), ("Hardware", 2.0), ("Toolchain", 2.2)],
             [["R1", "ESP32-D0WDQ6", "ESP-IDF 6.0.2"],
              ["R2", "Raspberry Pi 4B", "GCC 14.2.0, Mbed TLS 3.6.5"]],
             rh=0.36)
end2 = dtable(s, 7.2, end + 0.22, 5.5,
              [("Cycle", 1.0), ("R1 → R2 seq", 1.6), ("R2 → R1 seq", 1.6), ("ACK", 1.0)],
              [["1", "102762", "32769", "TX_DS"],
               ["2", "102763", "32770", "TX_DS"],
               ["3", "102764", "32771", "TX_DS"],
               ["4", "102765", "32772", "TX_DS"],
               ["5", "102766", "32773", "TX_DS"]],
              rh=0.3, size=10.5)
txt(s, 7.2, end2 + 0.2, 5.5, 0.62,
    "Sequence numbers strictly increasing on both nodes, which is the condition the anti-replay "
    "window relies on. Every frame 32 B, authenticated and decoded at the far end.",
    size=11.5, color=BODY, ls=1.15)
card(s, 7.2, end2 + 1.05, 5.5, 1.1, PALE)
txt(s, 7.45, end2 + 1.16, 5.0, 0.9,
    "Transport on this bench is Wi-Fi/UDP over an access point raised by the ESP32; the nRF24 "
    "FIFO, interrupt and retransmission behaviour is modelled in software.",
    size=11.5, color=BODY, anchor="middle", ls=1.15)
notes(s, """[1:45 - cumulative 18:45]

The protocol implementation running on physical processors.

Two nodes. An ESP32-D0WDQ6 as vehicle R1, firmware built with ESP-IDF 6.0.2 and flashed over a CP210x serial link. A Raspberry Pi 4 Model B as R2, with the native executables rebuilt on the Pi itself using GCC 14.2 and Mbed TLS 3.6.5. The same portable C source compiles for both.

The second table is the actual run from 16 August: five complete bidirectional cycles. R1 sends sequences 102762 through 102766, R2 answers with 32769 through 32773, and each transmission is acknowledged by the nRF24 model with a TX_DS.

What matters in that table is not the count but the monotonicity: sequence numbers strictly increasing on both nodes is exactly the condition the anti-replay window relies on. Every frame is 32 bytes, authenticated and decoded at the far end.

A detail worth mentioning: a sixth transmission from R1 ended in MAX_RT because R2 had exited after the five requested cycles. That is the missing-acknowledgement path being exercised, unintentionally but usefully.

The scope of the bench, on the card at the bottom: the transport is Wi-Fi and UDP over an access point raised by the ESP32, and the nRF24 FIFO, interrupt and retransmission behaviour is modelled in software. So what the bench establishes is the protocol, the cryptographic operations and the bidirectional scheduling on real silicon.""")

# =====================================================================
# 17 — SCHEDULE AND NEXT STEPS
# =====================================================================
s = slide()
kicker(s, "Schedule")
title(s, "Three months, and the bring-up sequence")
tl_y = 2.1
shape(s, MSO_SHAPE.RECTANGLE, M, tl_y, CW, 0.035, ICE)
months = [("June", "STM32/IMU bench\nSchematic and PCB corrected\nProduction file, order 23rd", DEEP),
          ("July", "Fleet simulator\nWeb supervision\nCapture scenario", TEAL),
          ("August", "HIL on two processors\nInductor re-selection\nReport", WARNC)]
mw = CW / 3.0
for i, (name, body, col) in enumerate(months):
    cx = M + i * mw + mw / 2
    shape(s, MSO_SHAPE.OVAL, cx - 0.11, tl_y - 0.09, 0.22, 0.22, col)
    txt(s, cx - mw / 2, tl_y - 0.72, mw, 0.4, name, size=21, bold=True, color=col,
        font=HEAD, align="center", anchor="middle")
    txt(s, cx - mw / 2 + 0.3, tl_y + 0.28, mw - 0.6, 1.1, body, size=12.5, color=BODY,
        align="center", ls=1.25)
card(s, M, 4.05, 12.1, 2.5, SURFACE)
txt(s, M + 0.4, 4.22, 7.0, 0.34, "Bring-up sequence on delivery", size=15, bold=True, color=DEEP)
nsteps = [("Continuity on +3V3 / GND before any power is applied", NOC, True),
          ("Current-limited ramp, measure both rails under load", BODY, False),
          ("Instrument the 5 V rail current — L1 has the tighter margin", BODY, False),
          ("Flash the STM32, bring peripherals up one bus at a time", BODY, False)]
for i, (t, col, hot) in enumerate(nsteps):
    x = M + 0.4 + (i % 2) * 5.85
    y = 4.7 + (i // 2) * 0.78
    badge(s, x, y, 0.44, "alert" if hot else "check", NOC if hot else DEEP)
    txt(s, x + 0.62, y - 0.03, 5.0, 0.5, t, size=12.5, bold=hot, color=col,
        anchor="middle", ls=1.06)
notes(s, """[1:00 - cumulative 19:45]

The distribution of the three months.

June: the exploratory STM32 and inertial bench, the schematic and PCB corrections, the production file and the order on the 23rd. July: the fleet simulator, the web supervision interface and the capture scenario. August: the HIL on two processors, the inductor re-selection after the PCBA review, and the report.

What carries beyond the internship: the simulator and its test suite specify the expected fleet behaviour independently of any particular board, so they serve as a reference for the final firmware. The protocol implementation is portable C already running on two processors. And the HIL bench remains valid when the Wi-Fi transport is replaced by the nRF24 link.

The bring-up sequence is the order I would follow. Continuity on the 3.3 volt and ground rails before any power is applied. Then a current-limited ramp with both rails measured under load. Then instrument the 5 volt rail current specifically, because L1 is the tighter margin at 8.7 percent. Then flash the STM32 and bring peripherals up one bus at a time.

One procedural point I am keeping: design rules, justified exclusions and the exported production files are versioned alongside the schematic, so the order is reproducible from the repository.""")

# =====================================================================
# 18 — CLOSING
# =====================================================================
s = slide(INK)
kicker(s, "Summary", color=ACCENT)
txt(s, M, 1.55, 12.1, 1.9,
    "A manufacturable production file for a fleet-capable AUV board,\n"
    "a deterministic mission model, and a portable protocol\n"
    "running authenticated on two processors.",
    size=27, bold=True, color=WHITE, font=HEAD, anchor="middle", ls=1.32)
st = [("94.8", "cm², 2 layers"), ("116", "components"), ("672", "W·h nominal"),
      ("12.3", "km survey track"), ("32", "B authenticated frame")]
sw2 = 2.3
for i, (v, l) in enumerate(st):
    x = M + i * (sw2 + 0.15)
    card(s, x, 4.05, sw2, 1.1, DARKCARD)
    txt(s, x + 0.08, 4.18, sw2 - 0.16, 0.44, v, size=23, bold=True,
        color=ACCENT, font=HEAD, align="center", anchor="middle")
    txt(s, x + 0.08, 4.63, sw2 - 0.16, 0.4, l, size=11, color=ICE, align="center")
txt(s, M, 5.6, 6.0, 0.44, "Thank you. I am available for questions.", size=18, bold=True,
    color=ACCENT, anchor="middle")
txt(s, M, 6.24, 8.0, 0.34,
    "Joab da Silva Bezerra  ·  CNRS – Lab-STICC  ·  ROBEX  ·  24 August 2026",
    size=11.5, color=MUTED, anchor="middle")
logos(s)
notes(s, """[0:30 - cumulative 20:15]

To summarise. The internship produced a manufacturable production file for a board capable of operating in a fleet, a deterministic mission model that specifies the fleet behaviour, and a portable protocol implementation running authenticated exchanges on two processors.

The figures on the strip are the ones that characterise the work: a 94.8 square centimetre two-layer board with 116 components, 672 watt-hours of nominal energy, 12.3 kilometres of survey track for the reference mission, and a 32-byte authenticated frame.

Thank you. I am available for questions.

Then stop talking.""")

prs.save(OUT)
print("WROTE", OUT)
print("slides:", len(prs.slides._sldIdLst))
